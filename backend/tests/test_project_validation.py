import tempfile
import unittest
from pathlib import Path

from config import settings
from services.project_validation import detect_project_profile, validate_project_after_code_agent


class FakeBus:
    def __init__(self) -> None:
        self.events = []

    async def publish(self, task_id, event_type, payload):
        self.events.append({"task_id": task_id, "type": event_type, "payload": payload})


class ProjectValidationTests(unittest.IsolatedAsyncioTestCase):
    def test_detects_python_project_profile(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            project = Path(tmp)
            (project / "main.py").write_text("print('ok')\n", encoding="utf-8")

            profile = detect_project_profile(project)

        self.assertEqual(profile["kind"], "python")
        self.assertTrue(profile["has_code"])

    async def test_python_project_passes_compile_validation(self) -> None:
        previous_workspace = settings.WORKSPACE_PATH
        with tempfile.TemporaryDirectory() as tmp:
            try:
                settings.WORKSPACE_PATH = Path(tmp)
                task_dir = settings.WORKSPACE_PATH / "task-python-ok"
                task_dir.mkdir()
                (task_dir / "main.py").write_text("def add(a, b):\n    return a + b\n", encoding="utf-8")
                (task_dir / "DOCUMENTACAO.md").write_text("# Script Python\n\n" + "Documentacao completa do script. " * 8, encoding="utf-8")
                bus = FakeBus()

                result = await validate_project_after_code_agent(
                    "task-python-ok",
                    "openclaude 'crie um script python'",
                    bus,
                    agent_result={"success": True},
                )
            finally:
                settings.WORKSPACE_PATH = previous_workspace

        self.assertEqual(result["status"], "passed")
        self.assertTrue(any(event["type"] == "project_validation_result" for event in bus.events))

    async def test_python_project_fails_compile_validation(self) -> None:
        previous_workspace = settings.WORKSPACE_PATH
        with tempfile.TemporaryDirectory() as tmp:
            try:
                settings.WORKSPACE_PATH = Path(tmp)
                task_dir = settings.WORKSPACE_PATH / "task-python-bug"
                task_dir.mkdir()
                (task_dir / "main.py").write_text("def broken(:\n    pass\n", encoding="utf-8")
                (task_dir / "DOCUMENTACAO.md").write_text("# Script Python\n\n" + "Documentacao completa do script. " * 8, encoding="utf-8")
                bus = FakeBus()

                result = await validate_project_after_code_agent(
                    "task-python-bug",
                    "openclaude 'crie um script python'",
                    bus,
                    agent_result={"success": True},
                )
            finally:
                settings.WORKSPACE_PATH = previous_workspace

        self.assertEqual(result["status"], "failed")
        self.assertTrue(any("py_compile" in bug for bug in result["bugs"]))

    async def test_static_project_fails_missing_asset_validation(self) -> None:
        previous_workspace = settings.WORKSPACE_PATH
        with tempfile.TemporaryDirectory() as tmp:
            try:
                settings.WORKSPACE_PATH = Path(tmp)
                task_dir = settings.WORKSPACE_PATH / "task-site-bug"
                task_dir.mkdir()
                (task_dir / "index.html").write_text('<script src="missing.js"></script>', encoding="utf-8")
                (task_dir / "DOCUMENTACAO.md").write_text("# Site\n\n" + "Documentacao completa do site. " * 8, encoding="utf-8")
                bus = FakeBus()

                result = await validate_project_after_code_agent(
                    "task-site-bug",
                    "openclaude 'crie um site html'",
                    bus,
                    agent_result={"success": True},
                )
            finally:
                settings.WORKSPACE_PATH = previous_workspace

        self.assertEqual(result["status"], "failed")
        self.assertTrue(any("missing.js" in bug for bug in result["bugs"]))

    async def test_site_project_requires_markdown_documentation(self) -> None:
        previous_workspace = settings.WORKSPACE_PATH
        with tempfile.TemporaryDirectory() as tmp:
            try:
                settings.WORKSPACE_PATH = Path(tmp)
                task_dir = settings.WORKSPACE_PATH / "task-site-no-doc"
                task_dir.mkdir()
                (task_dir / "index.html").write_text("<!doctype html><title>Site</title>", encoding="utf-8")
                bus = FakeBus()

                result = await validate_project_after_code_agent(
                    "task-site-no-doc",
                    "openclaude 'crie um site html'",
                    bus,
                    agent_result={"success": True},
                )
            finally:
                settings.WORKSPACE_PATH = previous_workspace

        self.assertEqual(result["status"], "failed")
        self.assertTrue(any("DOCUMENTACAO.md" in bug for bug in result["bugs"]))

    async def test_python_project_requires_markdown_documentation(self) -> None:
        previous_workspace = settings.WORKSPACE_PATH
        with tempfile.TemporaryDirectory() as tmp:
            try:
                settings.WORKSPACE_PATH = Path(tmp)
                task_dir = settings.WORKSPACE_PATH / "task-python-no-doc"
                task_dir.mkdir()
                (task_dir / "main.py").write_text("print('ok')\n", encoding="utf-8")
                bus = FakeBus()

                result = await validate_project_after_code_agent(
                    "task-python-no-doc",
                    "openclaude 'crie um script python'",
                    bus,
                    agent_result={"success": True},
                )
            finally:
                settings.WORKSPACE_PATH = previous_workspace

        self.assertEqual(result["status"], "failed")
        self.assertTrue(any("DOCUMENTACAO.md" in bug for bug in result["bugs"]))

    async def test_site_project_passes_with_markdown_documentation(self) -> None:
        previous_workspace = settings.WORKSPACE_PATH
        with tempfile.TemporaryDirectory() as tmp:
            try:
                settings.WORKSPACE_PATH = Path(tmp)
                task_dir = settings.WORKSPACE_PATH / "task-site-doc"
                task_dir.mkdir()
                (task_dir / "index.html").write_text("<!doctype html><title>Site</title>", encoding="utf-8")
                (task_dir / "DOCUMENTACAO.md").write_text("# Site\n\n" + "Documentacao completa do projeto. " * 8, encoding="utf-8")
                bus = FakeBus()

                result = await validate_project_after_code_agent(
                    "task-site-doc",
                    "openclaude 'crie um site html'",
                    bus,
                    agent_result={"success": True},
                )
            finally:
                settings.WORKSPACE_PATH = previous_workspace

        self.assertEqual(result["status"], "passed")

    async def test_document_request_requires_requested_extension(self) -> None:
        previous_workspace = settings.WORKSPACE_PATH
        with tempfile.TemporaryDirectory() as tmp:
            try:
                settings.WORKSPACE_PATH = Path(tmp)
                task_dir = settings.WORKSPACE_PATH / "task-pdf"
                task_dir.mkdir()
                (task_dir / "relatorio.md").write_text("# Relatorio\n\n" + "Conteudo completo do relatorio. " * 8, encoding="utf-8")
                bus = FakeBus()

                result = await validate_project_after_code_agent(
                    "task-pdf",
                    "openclaude 'gere um relatorio em PDF'",
                    bus,
                    agent_result={"success": True},
                )
            finally:
                settings.WORKSPACE_PATH = previous_workspace

        self.assertEqual(result["status"], "failed")
        self.assertTrue(any(".pdf" in bug for bug in result["bugs"]))

    async def test_pdf_request_passes_with_valid_pdf_and_markdown_source(self) -> None:
        previous_workspace = settings.WORKSPACE_PATH
        with tempfile.TemporaryDirectory() as tmp:
            try:
                settings.WORKSPACE_PATH = Path(tmp)
                task_dir = settings.WORKSPACE_PATH / "task-pdf-ok"
                task_dir.mkdir()
                (task_dir / "relatorio.md").write_text("# Relatorio\n\n" + "Conteudo pesquisado. " * 12, encoding="utf-8")
                (task_dir / "relatorio.pdf").write_bytes(b"%PDF-1.4\n" + b"x" * 300)
                bus = FakeBus()

                result = await validate_project_after_code_agent(
                    "task-pdf-ok",
                    "openclaude 'gere um relatorio em PDF'",
                    bus,
                    agent_result={"success": True},
                )
            finally:
                settings.WORKSPACE_PATH = previous_workspace

        self.assertEqual(result["status"], "passed")

    async def test_office_documents_pass_real_file_validation(self) -> None:
        from docx import Document
        from openpyxl import Workbook
        from pptx import Presentation

        previous_workspace = settings.WORKSPACE_PATH
        with tempfile.TemporaryDirectory() as tmp:
            try:
                settings.WORKSPACE_PATH = Path(tmp)
                task_dir = settings.WORKSPACE_PATH / "task-office-ok"
                task_dir.mkdir()

                document = Document()
                document.add_heading("Relatorio Executivo", level=1)
                document.add_paragraph("Analise completa com contexto, detalhes, resultados e recomendacoes. " * 2)
                document.save(task_dir / "relatorio.docx")

                presentation = Presentation()
                slide = presentation.slides.add_slide(presentation.slide_layouts[0])
                slide.shapes.title.text = "Resumo Executivo"
                slide.placeholders[1].text = "Contexto, dados principais, recomendacoes e proximos passos."
                presentation.save(task_dir / "slides.pptx")

                workbook = Workbook()
                sheet = workbook.active
                sheet.append(["Indicador", "Valor"])
                sheet.append(["Receita", 1200])
                sheet.append(["Custo", 800])
                workbook.save(task_dir / "dados.xlsx")

                (task_dir / "dados.csv").write_text("indicador,valor\nreceita,1200\ncusto,800\n", encoding="utf-8")
                bus = FakeBus()

                result = await validate_project_after_code_agent(
                    "task-office-ok",
                    "openclaude 'gere um docx, slides pptx, excel xlsx e csv'",
                    bus,
                    agent_result={"success": True},
                )
            finally:
                settings.WORKSPACE_PATH = previous_workspace

        self.assertEqual(result["status"], "passed")

    async def test_office_documents_fail_when_corrupted(self) -> None:
        previous_workspace = settings.WORKSPACE_PATH
        with tempfile.TemporaryDirectory() as tmp:
            try:
                settings.WORKSPACE_PATH = Path(tmp)
                task_dir = settings.WORKSPACE_PATH / "task-office-bug"
                task_dir.mkdir()
                (task_dir / "relatorio.docx").write_text("nao e docx real", encoding="utf-8")
                bus = FakeBus()

                result = await validate_project_after_code_agent(
                    "task-office-bug",
                    "openclaude 'gere um arquivo docx'",
                    bus,
                    agent_result={"success": True},
                )
            finally:
                settings.WORKSPACE_PATH = previous_workspace

        self.assertEqual(result["status"], "failed")
        self.assertTrue(any("DOCX" in bug for bug in result["bugs"]))

    async def test_github_repo_analysis_validates_only_technical_report(self) -> None:
        previous_workspace = settings.WORKSPACE_PATH
        with tempfile.TemporaryDirectory() as tmp:
            try:
                settings.WORKSPACE_PATH = Path(tmp)
                task_dir = settings.WORKSPACE_PATH / "task-github-analysis"
                repo_dir = task_dir / "requests"
                repo_dir.mkdir(parents=True)
                (repo_dir / "broken.py").write_text("def broken(:\n    pass\n", encoding="utf-8")
                (repo_dir / "package.json").write_text('{"scripts":{"test":"exit 1"}}', encoding="utf-8")
                (task_dir / "RELATORIO_TECNICO.md").write_text(
                    "# Analise do repositorio\n\n" + "Achado tecnico com evidencias por arquivo. " * 12,
                    encoding="utf-8",
                )
                bus = FakeBus()

                result = await validate_project_after_code_agent(
                    "task-github-analysis",
                    "openclaude 'analise https://github.com/psf/requests'",
                    bus,
                    agent_result={"success": True},
                )
            finally:
                settings.WORKSPACE_PATH = previous_workspace

        self.assertEqual(result["status"], "passed")
        self.assertEqual(result["checks"][0]["command"], "github repository analysis report scan")
        self.assertFalse(any("py_compile" in check["command"] for check in result["checks"]))

    async def test_github_repo_analysis_requires_root_technical_report(self) -> None:
        previous_workspace = settings.WORKSPACE_PATH
        with tempfile.TemporaryDirectory() as tmp:
            try:
                settings.WORKSPACE_PATH = Path(tmp)
                task_dir = settings.WORKSPACE_PATH / "task-github-no-report"
                repo_dir = task_dir / "repo"
                repo_dir.mkdir(parents=True)
                (repo_dir / "README.md").write_text("# Repo\n\n" + "Readme existente do projeto. " * 12, encoding="utf-8")
                bus = FakeBus()

                result = await validate_project_after_code_agent(
                    "task-github-no-report",
                    "openclaude 'analise https://github.com/example/repo'",
                    bus,
                    agent_result={"success": True},
                )
            finally:
                settings.WORKSPACE_PATH = previous_workspace

        self.assertEqual(result["status"], "failed")
        self.assertTrue(any("RELATORIO_TECNICO.md" in bug for bug in result["bugs"]))


if __name__ == "__main__":
    unittest.main()
