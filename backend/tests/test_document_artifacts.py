import tempfile
import unittest
import zipfile
from pathlib import Path

from config import settings
from services.document_artifacts import (
    archive_edit_targets,
    artifact_profile,
    csv_file_valid,
    docx_file_valid,
    markdown_file_valid,
    markdown_to_html,
    pdf_file_valid,
    pptx_file_valid,
    resolve_document_target,
    shapefile_zip_valid,
    valid_document_files,
    xlsx_file_valid,
    zip_file_valid,
)
from tools.tool_executor import _requested_document_artifact_error


class DocumentArtifactsTests(unittest.TestCase):
    def _write_shapefile_zip(self, target: Path) -> None:
        import geopandas as gpd
        from shapely.geometry import Point

        source_dir = target.parent / "shape-src"
        source_dir.mkdir()
        gdf = gpd.GeoDataFrame(
            {"nome": ["A", "B"], "valor": [10, 20]},
            geometry=[Point(-46.63, -23.55), Point(-46.64, -23.56)],
            crs="EPSG:4326",
        )
        shp_path = source_dir / "pontos.shp"
        gdf.to_file(shp_path, driver="ESRI Shapefile")
        with zipfile.ZipFile(target, "w", zipfile.ZIP_DEFLATED) as archive:
            for path in source_dir.iterdir():
                if path.is_file():
                    archive.write(path, f"pontos/{path.name}")

    def test_pdf_request_prefers_markdown_source_and_pdf_names(self) -> None:
        profile = artifact_profile("Gere um PDF com a historia do Corinthians")

        self.assertTrue(profile["wants_pdf"])
        self.assertTrue(profile["requires_artifact"])
        self.assertTrue(profile["preferred_markdown"].endswith(".md"))
        self.assertTrue(profile["preferred_pdf"].endswith(".pdf"))
        self.assertIn("corinthians", profile["preferred_pdf"])

    def test_detects_contextual_document_edit(self) -> None:
        profile = artifact_profile("melhore esse PDF e deixe mais completo")

        self.assertTrue(profile["edit_requested"])
        self.assertTrue(profile["wants_pdf"])

    def test_resolves_latest_previewable_document_target(self) -> None:
        files = [
            {"path": "old.pdf", "extension": ".pdf", "modified_at": 1, "size_bytes": 300},
            {"path": "versions/old.pdf", "extension": ".pdf", "modified_at": 3, "size_bytes": 300},
            {"path": "novo.pdf", "extension": ".pdf", "modified_at": 2, "size_bytes": 300},
        ]

        target = resolve_document_target("task", "melhore o pdf", files)

        self.assertEqual(target["path"], "novo.pdf")

    def test_archives_existing_document_for_edit(self) -> None:
        previous_workspace = settings.WORKSPACE_PATH
        with tempfile.TemporaryDirectory() as tmp:
            try:
                settings.WORKSPACE_PATH = Path(tmp)
                task_dir = settings.WORKSPACE_PATH / "task-doc"
                task_dir.mkdir()
                (task_dir / "relatorio.pdf").write_bytes(b"%PDF-1.4\n" + b"x" * 300)
                archived = archive_edit_targets(
                    "task-doc",
                    "melhore o pdf",
                    [{"path": "relatorio.pdf", "extension": ".pdf", "modified_at": 1, "size_bytes": 309}],
                )
            finally:
                settings.WORKSPACE_PATH = previous_workspace

        self.assertEqual(len(archived), 1)
        self.assertTrue(archived[0].startswith("versions/relatorio_"))

    def test_validates_markdown_and_pdf_content(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            md = root / "doc.md"
            pdf = root / "doc.pdf"
            md.write_text("# Titulo\n\n" + "conteudo " * 20, encoding="utf-8")
            pdf.write_bytes(b"%PDF-1.4\n" + b"x" * 300)

            self.assertTrue(markdown_file_valid(md))
            self.assertTrue(pdf_file_valid(pdf))

    def test_validates_office_and_csv_content(self) -> None:
        from docx import Document
        from openpyxl import Workbook
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)

            docx_path = root / "relatorio.docx"
            document = Document()
            document.add_heading("Relatorio Executivo", level=1)
            document.add_paragraph("Conteudo completo do documento com detalhes, contexto, analise e conclusoes. " * 2)
            document.save(docx_path)

            pptx_path = root / "apresentacao.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            slide.shapes.title.text = "Apresentacao Executiva"
            slide.placeholders[1].text = "Resumo, contexto, dados principais e proximos passos."
            presentation.save(pptx_path)

            xlsx_path = root / "planilha.xlsx"
            workbook = Workbook()
            sheet = workbook.active
            sheet.append(["Item", "Valor"])
            sheet.append(["Receita", 1200])
            sheet.append(["Custo", 800])
            workbook.save(xlsx_path)

            csv_path = root / "dados.csv"
            csv_path.write_text("item,valor\nreceita,1200\ncusto,800\n", encoding="utf-8")

            self.assertTrue(docx_file_valid(docx_path))
            self.assertTrue(pptx_file_valid(pptx_path))
            self.assertTrue(xlsx_file_valid(xlsx_path))
            self.assertTrue(csv_file_valid(csv_path))
            self.assertEqual(valid_document_files(root, ".docx"), ["relatorio.docx"])

    def test_validates_zip_and_shapefile_zip_content(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            generic_zip = root / "arquivos.zip"
            with zipfile.ZipFile(generic_zip, "w", zipfile.ZIP_DEFLATED) as archive:
                archive.writestr("relatorio.txt", "conteudo")

            shape_zip = root / "shape_final.zip"
            self._write_shapefile_zip(shape_zip)

            invalid_shape_zip = root / "shape_incompleto.zip"
            with zipfile.ZipFile(invalid_shape_zip, "w", zipfile.ZIP_DEFLATED) as archive:
                archive.writestr("layer.shp", b"fake")

            self.assertTrue(zip_file_valid(generic_zip))
            self.assertFalse(shapefile_zip_valid(generic_zip))
            self.assertTrue(shapefile_zip_valid(shape_zip))
            self.assertFalse(shapefile_zip_valid(invalid_shape_zip))
            self.assertEqual(valid_document_files(root, ".zip"), ["arquivos.zip", "shape_final.zip", "shape_incompleto.zip"])

    def test_markdown_to_html_preserves_headings(self) -> None:
        html = markdown_to_html("# Titulo\n\n## Secao\n\nTexto")

        self.assertIn("<h1>Titulo</h1>", html)
        self.assertIn("<h2>Secao</h2>", html)

    def test_pdf_request_requires_real_markdown_and_pdf_artifacts(self) -> None:
        previous_workspace = settings.WORKSPACE_PATH
        with tempfile.TemporaryDirectory() as tmp:
            try:
                settings.WORKSPACE_PATH = Path(tmp)
                task_dir = settings.WORKSPACE_PATH / "task-pdf"
                task_dir.mkdir()

                missing = _requested_document_artifact_error("task-pdf", "openclaude 'gere um relatorio em PDF'")
                (task_dir / "relatorio.md").write_text("# Relatorio\n\n" + "conteudo " * 20, encoding="utf-8")
                only_markdown = _requested_document_artifact_error("task-pdf", "openclaude 'gere um relatorio em PDF'")
                (task_dir / "relatorio.pdf").write_bytes(b"%PDF-1.4\n" + b"x" * 300)
                complete = _requested_document_artifact_error("task-pdf", "openclaude 'gere um relatorio em PDF'")
            finally:
                settings.WORKSPACE_PATH = previous_workspace

        self.assertIn("Markdown fonte", missing)
        self.assertIn(".pdf valido", only_markdown)
        self.assertIsNone(complete)

    def test_office_request_requires_real_artifact(self) -> None:
        previous_workspace = settings.WORKSPACE_PATH
        with tempfile.TemporaryDirectory() as tmp:
            try:
                settings.WORKSPACE_PATH = Path(tmp)
                task_dir = settings.WORKSPACE_PATH / "task-docx"
                task_dir.mkdir()

                missing = _requested_document_artifact_error("task-docx", "openclaude 'gere um arquivo docx'")
                (task_dir / "relatorio.docx").write_text("isto nao e um docx real", encoding="utf-8")
                invalid = _requested_document_artifact_error("task-docx", "openclaude 'gere um arquivo docx'")

                from docx import Document

                document = Document()
                document.add_heading("Relatorio", level=1)
                document.add_paragraph("Conteudo estruturado, completo e suficiente para validar o arquivo DOCX. " * 2)
                document.save(task_dir / "relatorio.docx")
                complete = _requested_document_artifact_error("task-docx", "openclaude 'gere um arquivo docx'")
            finally:
                settings.WORKSPACE_PATH = previous_workspace

        self.assertIn(".docx real", missing)
        self.assertIn(".docx real", invalid)
        self.assertIsNone(complete)

    def test_shapefile_edit_request_requires_final_zip_with_complete_shape(self) -> None:
        previous_workspace = settings.WORKSPACE_PATH
        with tempfile.TemporaryDirectory() as tmp:
            try:
                settings.WORKSPACE_PATH = Path(tmp)
                task_dir = settings.WORKSPACE_PATH / "task-shape"
                task_dir.mkdir()

                missing = _requested_document_artifact_error(
                    "task-shape",
                    "vertex 'edite a tabela de atributos do shapefile'",
                )
                bad_zip = task_dir / "outputs.zip"
                with zipfile.ZipFile(bad_zip, "w", zipfile.ZIP_DEFLATED) as archive:
                    archive.writestr("layer.shp", b"fake")
                invalid = _requested_document_artifact_error(
                    "task-shape",
                    "vertex 'edite a tabela de atributos do shapefile'",
                )

                good_zip = task_dir / "shape_editado.zip"
                self._write_shapefile_zip(good_zip)
                complete = _requested_document_artifact_error(
                    "task-shape",
                    "vertex 'edite a tabela de atributos do shapefile'",
                )
            finally:
                settings.WORKSPACE_PATH = previous_workspace

        self.assertIn(".zip final", missing)
        self.assertIn(".zip final", invalid)
        self.assertIsNone(complete)


if __name__ == "__main__":
    unittest.main()
